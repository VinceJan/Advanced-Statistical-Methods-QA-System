from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables"
SHOT = ROOT / "docs" / "qa-screenshots"

W = Inches(13.333)
H = Inches(7.5)

INK = RGBColor(23, 32, 51)
MUTED = RGBColor(102, 117, 138)
BG = RGBColor(237, 244, 247)
PANEL = RGBColor(255, 255, 255)
TEAL = RGBColor(22, 113, 139)
TEAL_DARK = RGBColor(15, 42, 58)
AMBER = RGBColor(159, 98, 56)
GREEN = RGBColor(19, 121, 91)
LINE = RGBColor(216, 225, 234)
SOFT = RGBColor(243, 247, 250)
WARN = RGBColor(173, 106, 0)


def add_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=24, color=INK, bold=False, align=None, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    p = frame.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.72, 0.45, 9.8, 0.55, 25, TEAL_DARK, True)
    if subtitle:
        add_text(slide, subtitle, 0.74, 1.05, 10.5, 0.32, 10.5, MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(1.48), Inches(11.9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_card(slide, x, y, w, h, title=None, body=None, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    if title:
        add_text(slide, title, x + 0.18, y + 0.16, w - 0.35, 0.32, 13, accent, True)
    if body:
        add_text(slide, body, x + 0.18, y + 0.55, w - 0.35, h - 0.68, 9.5, INK)
    return shape


def add_metric(slide, x, y, label, value, color=TEAL):
    add_card(slide, x, y, 1.75, 0.9)
    add_text(slide, str(value), x + 0.18, y + 0.14, 1.25, 0.32, 20, color, True)
    add_text(slide, label, x + 0.2, y + 0.55, 1.25, 0.22, 8.5, MUTED)


def add_bullets(slide, items, x, y, w, h, size=13, color=INK, gap=0.33):
    top = y
    for item in items:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(top + 0.08), Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL
        dot.line.fill.background()
        add_text(slide, item, x + 0.18, top, w - 0.18, 0.34, size, color)
        top += gap


def add_image(slide, path, x, y, w, h):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        add_card(slide, x, y, w, h, "截图缺失", str(path), WARN)


def add_footer(slide, idx, total):
    add_text(slide, "高级统计方法智能问答系统", 0.72, 7.08, 4.2, 0.2, 7.5, MUTED)
    add_text(slide, f"{idx:02d}/{total:02d}", 11.75, 7.08, 0.8, 0.2, 7.5, MUTED)


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    return slide


def cover(prs, title, subtitle, label):
    slide = blank(prs)
    add_bg(slide, TEAL_DARK)
    add_text(slide, label, 0.78, 0.62, 4.0, 0.25, 9.5, RGBColor(196, 222, 229), True)
    add_text(slide, title, 0.78, 1.28, 8.8, 1.05, 32, RGBColor(255, 255, 255), True)
    add_text(slide, subtitle, 0.82, 2.55, 7.7, 0.52, 14, RGBColor(217, 240, 242))
    add_metric(slide, 0.82, 4.55, "知识点", 69, RGBColor(217, 240, 242))
    add_metric(slide, 2.82, 4.55, "图谱边", 336, RGBColor(217, 240, 242))
    add_metric(slide, 4.82, 4.55, "问答对", 70, RGBColor(217, 240, 242))
    add_metric(slide, 6.82, 4.55, "文本块", 1413, RGBColor(217, 240, 242))
    add_image(slide, SHOT / "remote-graph.png", 8.65, 1.0, 4.0, 3.05)
    add_text(slide, "已部署：https://aistudyassistant.bluesclawd.dev", 0.82, 6.55, 5.5, 0.25, 10, RGBColor(217, 240, 242))


def section_slide(prs, title, points, idx, total):
    slide = blank(prs)
    add_text(slide, title, 0.8, 1.0, 8.8, 0.7, 30, TEAL_DARK, True)
    add_bullets(slide, points, 0.95, 2.15, 7.8, 2.8, 15, INK, 0.48)
    add_card(slide, 9.3, 1.15, 3.0, 4.85, "本节关键词", "RAG 可信度\n知识图谱\n工程验证\n公网部署", AMBER)
    add_footer(slide, idx, total)


def notes_md(title, slides):
    lines = [f"# {title}", ""]
    for i, s in enumerate(slides, 1):
        lines += [f"## {i}. {s['title']}", "", s["note"], ""]
    return "\n".join(lines)


def finalize(prs, path):
    prs.save(path)


def deck_base():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def slide_simple(prs, title, subtitle, bullets, image=None, metrics=None):
    slide = blank(prs)
    add_title(slide, title, subtitle)
    if metrics:
        x = 0.78
        for label, value, color in metrics:
            add_metric(slide, x, 1.85, label, value, color)
            x += 2.0
        y = 3.05
    else:
        y = 1.95
    if image:
        add_bullets(slide, bullets, 0.85, y, 5.15, 4.2, 12.5, INK, 0.42)
        add_image(slide, SHOT / image, 6.25, 1.85, 6.1, 4.25)
    else:
        add_bullets(slide, bullets, 0.88, y, 10.7, 4.8, 14, INK, 0.47)
    return slide


def slide_architecture(prs):
    slide = blank(prs)
    add_title(slide, "系统架构：从课程问题到可信回答", "FastAPI + React + SQLite + 教材 RAG + 知识图谱 + MiniMax")
    labels = [
        ("用户问题", 0.8, 2.1, TEAL),
        ("范围判断", 2.8, 2.1, AMBER),
        ("教材检索", 4.8, 2.1, TEAL),
        ("图谱子图", 6.8, 2.1, GREEN),
        ("MiniMax 生成", 8.8, 2.1, TEAL),
        ("中文答案", 10.8, 2.1, AMBER),
    ]
    for text, x, y, color in labels:
        add_card(slide, x, y, 1.55, 0.72, text, None, color)
        if x < 10:
            add_text(slide, "→", x + 1.63, y + 0.18, 0.28, 0.25, 18, MUTED, True)
    add_card(slide, 0.85, 3.45, 3.55, 1.55, "可信度守门", "out_of_scope / insufficient_evidence\n低证据不展示伪来源和伪图谱", AMBER)
    add_card(slide, 4.75, 3.45, 3.55, 1.55, "数据闭环", "会话、历史、用户、问答对、文本块、知识点和图谱边全部结构化入库", TEAL)
    add_card(slide, 8.65, 3.45, 3.55, 1.55, "管理闭环", "管理员可配置模型、维护数据资产、审计历史并管理学生账号", GREEN)


def slide_requirements(prs):
    slide = blank(prs)
    add_title(slide, "题目要求映射：不是 demo，而是完整交付", "围绕数据规模、功能闭环、验证材料和课程报告逐项覆盖")
    rows = [
        ("数据资产", "69 知识点 / 336 图谱边 / 70 问答对 / 1413 文本块", GREEN),
        ("问答功能", "中文多轮追问、RAG 来源、图谱子图、推荐问题、历史全文", TEAL),
        ("管理功能", "用户权限、模型 Key、问答对、知识点、图谱边、文本块和历史审计", AMBER),
        ("验证材料", "pytest、构建、npm audit、Playwright 截图、公网验收", GREEN),
        ("文档报告", "README、AGENTS、部署方案、对抗性审查、实验报告、PPT", TEAL),
    ]
    y = 1.85
    for left, right, color in rows:
        add_card(slide, 0.9, y, 2.3, 0.72, left, None, color)
        add_card(slide, 3.45, y, 8.6, 0.72, right, None, color)
        y += 0.88


def slide_deployment(prs):
    slide = blank(prs)
    add_title(slide, "部署方案：Docker Compose + 宿主机 Caddy", "公网域名已接入，验证覆盖健康检查、登录、问答和知识图谱")
    add_card(slide, 0.85, 1.85, 3.0, 1.3, "公网入口", "https://aistudyassistant.bluesclawd.dev\nCaddy 自动 HTTPS", TEAL)
    add_card(slide, 4.05, 1.85, 3.0, 1.3, "容器服务", "backend: 127.0.0.1:18000\nfrontend: 127.0.0.1:18080", GREEN)
    add_card(slide, 7.25, 1.85, 3.0, 1.3, "持久化", "Docker volume 保存 SQLite 数据\n.env 保存密钥和管理员密码", AMBER)
    add_card(slide, 10.45, 1.85, 1.7, 1.3, "仓库", "GitHub public\nmain", TEAL)
    add_image(slide, SHOT / "remote-graph.png", 0.95, 3.55, 5.7, 2.85)
    add_bullets(slide, [
        "线上健康检查：/api/health 成功",
        "线上统计：69 / 336 / 70 / 1413",
        "真实课程问答：answer_mode=minimax",
        "默认管理员密码已在服务器侧替换",
    ], 7.05, 3.72, 4.75, 2.2, 12.5, INK, 0.42)


def build_detailed():
    prs = deck_base()
    slides = []
    cover(prs, "高级统计方法知识图谱智能问答系统", "课程设计答辩详尽版 · RAG 可信问答 + 知识图谱 + 管理后台 + 公网部署", "Course Project Defense")
    slides.append({"title": "封面", "note": "开场说明本项目不是单纯聊天机器人，而是围绕高级统计方法课程内容构建的可信问答系统，已经完成本地和公网部署。"})
    data = [
        ("项目目标", "从题目二出发，将教材、知识图谱和问答对组织成可交互的课程学习工具。", ["解决课程知识分散、教材英文抽取、问答缺少引用、图谱难以交互的问题", "系统必须会拒答无关问题，不为了显得聪明而伪造来源", "最终交付包含代码、数据、验证截图、部署和汇报材料"], None),
        ("教材选择", "英文 ISLRv2 作为 RAG 主语料，中文扫描教材仅作术语参考。", ["中文教材是扫描版，OCR 成本高且噪声大", "英文 PDF 是文本型，能稳定抽取章节、页码和正文", "中文提问通过术语扩展检索英文证据，回答仍保持中文体验"], None),
        ("数据资产", "系统启动后自动种子化课程知识资产并建立教材索引。", ["知识点覆盖监督学习、回归、分类、重采样、正则化、树模型、SVM、无监督学习等", "图谱边包含“属于、相关、对比、用于、评估”等关系", "文本块经过清洗，过滤目录、图注、坐标轴和公式噪声"], None),
    ]
    for title, sub, bullets, image in data:
        slide_simple(prs, title, sub, bullets, image, [("知识点", 69, GREEN), ("图谱边", 336, TEAL), ("问答对", 70, AMBER), ("文本块", 1413, GREEN)] if title == "数据资产" else None)
        slides.append({"title": title, "note": "本页说明：" + "；".join(bullets)})
    slide_architecture(prs)
    slides.append({"title": "系统架构", "note": "讲解一次请求从前端进入后端，先做范围判断，再检索教材证据和图谱事实，最后调用 MiniMax 或本地降级回答。"})
    more = [
        ("RAG 可信度设计", "回答前先判定问题是否在课程范围内。", ["无关问题返回 out_of_scope，来源、推荐问题和图谱全部为空", "统计相关但教材证据不足时返回 insufficient_evidence", "正常问题才展示来源、相关问题和图谱子图"], "desktop-out-of-scope.png"),
        ("MiniMax 接入", "通过 OpenAI-compatible Chat Completions 调用 MiniMax-M3。", ["API Key 只在 .env 和运行时配置表保存", "管理员后台可切换 Base URL、模型和 Key", "自动化测试通过 APP_DISABLE_LLM=true 避免费用消耗"], "desktop-admin.png"),
        ("多轮追问", "会话消息结构化保存，每轮证据独立刷新。", ["追问会结合最近上下文补全问题", "清晰无关问题不会被历史上下文污染", "侧栏来源、推荐问题和图谱跟随当前轮次切换"], "desktop-follow-up.png"),
        ("学习历史", "历史记录不只是一行摘要，点击可展开当时全文。", ["每条历史保存回答、来源、conversation_id 和 message_id", "可用于课后复盘，也便于教师审计问答质量", "清空历史只影响当前用户"], "desktop-history-detail.png"),
        ("知识图谱交互", "Cytoscape 自动布局解决标签重叠和节点排布问题。", ["学生端可搜索概念并加载子图", "节点点击联动邻域和推荐问题", "英文术语完整显示，不再被按钮高度截断"], "desktop-student-graph.png"),
        ("管理员后台", "从内容维护扩展到系统运营。", ["管理用户角色、改密、注销会话和删除学生账号", "维护问答对、知识点、图谱边和教材文本块", "配置 LLM API 并审计历史提问"], "desktop-admin.png"),
    ]
    for title, sub, bullets, image in more:
        slide_simple(prs, title, sub, bullets, image)
        slides.append({"title": title, "note": "本页重点：" + "；".join(bullets)})
    slide_requirements(prs)
    slides.append({"title": "题目要求映射", "note": "这一页把系统实现和课程题目要求逐项对应，证明交付不是散点功能，而是一个完整闭环。"})
    detailed_only = [
        ("性能观测", "每次回答返回耗时拆分，定位慢点主要在 LLM 生成。", ["analysis_ms、retrieval_ms、graph_ms、llm_ms、total_ms 都会展示", "文本块索引使用进程内缓存，热查询检索开销较低", "公网验证中一次 MiniMax 回答约 19.8 秒，主要耗时来自生成"], None),
        ("自动化验证", "用脚本把“能跑”提升为“可验收”。", ["pytest 覆盖认证、问答、图谱、历史和管理权限", "npm build 与 npm audit 检查前端质量和依赖风险", "Playwright 生成桌面、移动端和远程部署截图"], None),
        ("对抗性审查", "从第一性原理检查系统是否会装懂。", ["无关问题、提示注入、要求编造来源、术语混淆、公式噪声都进入对抗题集", "默认图谱不再回退到固定中心节点", "API Key 永远不以明文出现在响应或文档中"], None),
    ]
    for title, sub, bullets, image in detailed_only:
        slide_simple(prs, title, sub, bullets, image)
        slides.append({"title": title, "note": "本页可以强调工程质量：" + "；".join(bullets)})
    slide_deployment(prs)
    slides.append({"title": "部署方案", "note": "介绍线上部署方式，说明为什么选择 Docker Compose 加宿主机 Caddy，以及如何验证公网链路。"})
    tail = [
        ("技术取舍", "保守地沿用轻量技术栈，优先完成课程项目闭环。", ["SQLite 足够支撑课程演示，未来可替换 PostgreSQL", "本地混合检索稳定可控，未来可升级 Chroma/FAISS", "Cytoscape 负责图谱可视化，避免手写 SVG 的排版风险"]),
        ("可继续优化", "系统已达到正式演示水平，但仍有产品化增长空间。", ["索引重建改为后台任务并显示进度", "增加 LLM token、失败率和成本统计", "批量导入导出问答对和知识图谱，方便教师验收"]),
        ("答辩总结", "本项目完成了从课程资料到可信智能助教的完整工程落地。", ["有数据资产、有交互系统、有管理后台、有验证链路、有公网部署", "重点不是让模型随便回答，而是让系统有边界、有依据、有可审查性", "可以作为课程设计作品直接展示和答辩"]),
    ]
    for title, sub, bullets in tail:
        slide_simple(prs, title, sub, bullets)
        slides.append({"title": title, "note": "收束陈述：" + "；".join(bullets)})
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        if i != 1:
            add_footer(slide, i, total)
    return prs, notes_md("高级统计方法智能问答系统答辩详尽版演讲稿", slides)


def build_medium():
    prs = deck_base()
    slides = []
    cover(prs, "高级统计方法智能问答系统", "课程项目汇报适中版 · 功能、架构、验证与部署", "Project Report")
    slides.append({"title": "封面", "note": "用一分钟概括项目定位：面向高级统计方法课程的中文可信问答系统。"})
    items = [
        ("项目定位", "用课程教材和知识图谱构建可信助教。", ["中文问答、多轮追问、来源引用、知识图谱和后台管理", "核心原则：不在范围内就拒答，证据不足不编造", "已完成本地与公网部署"]),
        ("数据与教材", "英文 ISLRv2 文本型 PDF 是主语料。", ["中文扫描教材不适合直接抽取", "英文教材稳定提供章节、页码和正文", "系统内置 69 知识点、336 边、70 问答对、1413 文本块"]),
        ("系统架构", "FastAPI + React + SQLite + MiniMax。", ["前端负责会话、来源、图谱和后台操作", "后端负责认证、RAG、图谱、历史和权限", "MiniMax 负责生成中文结构化回答"]),
        ("可信问答", "先判断范围，再检索证据。", ["无关问题 sources 和 graph 全为空", "低证据问题明确说明不足", "正常问题给中文回答、来源和图谱"],),
    ]
    for title, sub, bullets in items:
        slide_simple(prs, title, sub, bullets)
        slides.append({"title": title, "note": "讲解：" + "；".join(bullets)})
    slide_simple(prs, "功能展示：问答工作台", "支持连续追问，每轮独立刷新证据。", ["会话列表和消息流保存上下文", "回答区支持 Markdown 和 KaTeX", "来源、相关问题、图谱跟随当前轮次"], "desktop-follow-up.png")
    slides.append({"title": "问答工作台", "note": "展示多轮追问截图，强调证据不会沿用旧轮次。"})
    slide_simple(prs, "功能展示：知识图谱", "学生端和管理员端均可搜索和查看子图。", ["Cytoscape 自动布局", "英文术语完整显示", "点击节点联动推荐问题"], "desktop-student-graph.png")
    slides.append({"title": "知识图谱", "note": "展示图谱页面，说明修复了学生端加载和英文截断问题。"})
    slide_simple(prs, "管理后台", "管理员可以维护数据和运行配置。", ["用户角色、改密、注销、删除", "问答对、知识点、图谱边和文本块维护", "LLM API Key 与模型配置"], "desktop-admin.png")
    slides.append({"title": "管理后台", "note": "说明后台不是摆设，而是能真正维护系统运行状态。"})
    slide_requirements(prs)
    slides.append({"title": "题目要求映射", "note": "用一页把课程要求和系统交付对应起来。"})
    slide_deployment(prs)
    slides.append({"title": "部署与验证", "note": "说明公网部署和验证结果，给答辩演示留出可信支撑。"})
    slide_simple(prs, "总结与展望", "项目已经达到可正式演示的课程作品状态。", ["完成数据、系统、后台、验证、部署和汇报材料", "后续可升级向量数据库、后台任务和观测面板", "重点价值：可信、有边界、可审查、可演示"])
    slides.append({"title": "总结与展望", "note": "最后强调项目价值和未来扩展方向。"})
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        if i != 1:
            add_footer(slide, i, total)
    return prs, notes_md("高级统计方法智能问答系统汇报适中版演讲稿", slides)


def main():
    OUT.mkdir(exist_ok=True)
    detailed, detailed_notes = build_detailed()
    medium, medium_notes = build_medium()
    finalize(detailed, OUT / "高级统计方法智能问答系统_答辩详尽版.pptx")
    finalize(medium, OUT / "高级统计方法智能问答系统_汇报适中版.pptx")
    (OUT / "高级统计方法智能问答系统_答辩详尽版_演讲稿.md").write_text(detailed_notes, encoding="utf-8")
    (OUT / "高级统计方法智能问答系统_汇报适中版_演讲稿.md").write_text(medium_notes, encoding="utf-8")


if __name__ == "__main__":
    main()
